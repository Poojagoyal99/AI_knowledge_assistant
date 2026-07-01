import os
import shutil
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def _get_tesseract_cmd():
    if os.getenv("TESSERACT_CMD"):
        return os.getenv("TESSERACT_CMD")
    return shutil.which("tesseract")


def _ocr_pdf(pdf_path):
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        raise RuntimeError("OCR dependencies not installed")

    tesseract_cmd = _get_tesseract_cmd()
    if not tesseract_cmd:
        raise RuntimeError("Tesseract not found")

    pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

    images = convert_from_path(pdf_path, dpi=300)
    text = ""

    for image in images:
        page_text = pytesseract.image_to_string(image, lang="eng")
        if page_text:
            text += page_text + "\n"

    return text


def load_pdf(pdf_path):
    reader = PdfReader(pdf_path)
    text = ""

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text

    # ✅ If text exists → return immediately
    if text.strip():
        return text

    # ⚠️ Only try OCR if needed
    try:
        return _ocr_pdf(pdf_path)
    except Exception as e:
        print(f"OCR skipped: {e}")
        return ""


def load_docx(file_path):
    """Extract text from a Word .docx file."""
    from docx import Document

    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def load_pptx(file_path):
    """Extract text from a PowerPoint .pptx file."""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        slide_texts.append(para_text)
        if slide_texts:
            text_parts.append(f"[Slide {slide_num}] " + " ".join(slide_texts))
    return "\n".join(text_parts)


def load_txt(file_path):
    """Read a plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_document(file_path):
    """Load text from any supported file format."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".docx":
        return load_docx(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    elif ext == ".txt":
        return load_txt(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")