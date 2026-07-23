import os

from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def load_document(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _load_pdf(file_path)
    elif ext == ".docx":
        return _load_docx(file_path)
    elif ext == ".pptx":
        return _load_pptx(file_path)
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _load_pdf(path: str) -> str:
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def _load_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _load_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
    return "\n".join(text_parts)
